import os
from flask import Flask, request, jsonify
from flask_cors import CORS
from openai import OpenAI
import faiss, numpy as np
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
import requests  # 🌐 URL fetching

# -------------------------------------------------------------------
# ✅ APP SETUP
# -------------------------------------------------------------------
app = Flask(__name__)
CORS(app)

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# ---- Short-term memory (RAM only) ----
memory_vectors = []
memory_texts = []

# -------------------------------------------------------------------
# ✅ UTILITIES
# -------------------------------------------------------------------
def embed_text(text: str) -> np.ndarray:
    """Generate embedding vector for given text using OpenAI."""
    try:
        emb = client.embeddings.create(
            input=text,
            model="text-embedding-3-small"
        ).data[0].embedding
        return np.array(emb, dtype="float32")
    except Exception as e:
        print("Embedding error:", e)
        return np.zeros(1536, dtype="float32")


def extract_text(file):
    """Extract readable text from multiple document formats."""
    name = file.filename.lower()

    # PDF
    if name.endswith(".pdf"):
        try:
            reader = PdfReader(file)
            return "\n".join([p.extract_text() or "" for p in reader.pages])
        except Exception:
            return ""

    # DOCX
    elif name.endswith(".docx"):
        doc = Document(file)
        return "\n".join([p.text for p in doc.paragraphs])

    # EXCEL / CSV
    elif name.endswith((".xls", ".xlsx", ".csv")):
        try:
            df = pd.read_excel(file) if not name.endswith(".csv") else pd.read_csv(file)
            return df.to_string()
        except Exception:
            return ""

    # POWERPOINT
    elif name.endswith(".pptx"):
        try:
            prs = Presentation(file)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception:
            return ""

    # HTML
    elif name.endswith((".html", ".htm")):
        try:
            soup = BeautifulSoup(file.read(), "html.parser")
            return soup.get_text()
        except Exception:
            return ""

    # Plain text / others
    else:
        try:
            return file.read().decode("utf-8", errors="ignore")
        except Exception:
            return ""


def chunk_text(text, size=4000, overlap=200):
    """Split long text into safe-sized overlapping chunks."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + size
        chunks.append(text[start:end])
        start += size - overlap
    return chunks


# -------------------------------------------------------------------
# ✅ ROUTES
# -------------------------------------------------------------------
@app.route("/upload", methods=["POST"])
def upload_files():
    """Upload and embed files into short-term memory (chunked)."""
    global memory_vectors, memory_texts
    files = request.files.getlist("files")
    if not files:
        return jsonify({"error": "No files uploaded."}), 400

    count = 0
    for f in files:
        text = extract_text(f)
        if not text.strip():
            continue

        # Split long documents into manageable chunks
        for chunk in chunk_text(text):
            emb = embed_text(chunk)
            memory_vectors.append(emb)
            memory_texts.append(chunk)
        count += 1

    return jsonify({"message": f"✅ Uploaded {count} file(s) to short-term memory."})


# -------------------------------------------------------------------
# 🌐 NEW: Fetch and embed webpage content
# -------------------------------------------------------------------
@app.route("/url", methods=["POST"])
def upload_url():
    """Fetch and embed a webpage directly from a URL."""
    global memory_vectors, memory_texts
    data = request.get_json()
    url = data.get("url", "").strip()
    if not url:
        return jsonify({"error": "No URL provided."}), 400

    try:
        res = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
        res.raise_for_status()
        soup = BeautifulSoup(res.text, "html.parser")
        text = soup.get_text()
    except Exception as e:
        return jsonify({"error": f"Failed to fetch or parse: {e}"}), 500

    if not text.strip():
        return jsonify({"error": "No readable text found on the page."}), 400

    for chunk in chunk_text(text):
        emb = embed_text(chunk)
        memory_vectors.append(emb)
        memory_texts.append(chunk)

    return jsonify({"message": f"✅ Page from {url} added to short-term memory."})


# -------------------------------------------------------------------
# 🤖 ASK ENDPOINT WITH ANSWER LENGTH MODES
# -------------------------------------------------------------------
@app.route("/ask", methods=["POST"])
def ask_question():
    """Answer based on context with variable length modes."""
    global memory_vectors, memory_texts
    data = request.get_json()
    question = data.get("question", "").strip()
    length = data.get("length", "concise").lower()  # new parameter

    if not question:
        return jsonify({"error": "No question provided."}), 400
    if not memory_vectors:
        return jsonify({"answer": "Memory is empty. Please upload files or fetch a URL first."})

    # ---- Retrieve relevant context ----
    q_emb = embed_text(question)
    index = faiss.IndexFlatL2(len(q_emb))
    index.add(np.stack(memory_vectors))
    _, I = index.search(np.array([q_emb]), k=min(5, len(memory_vectors)))
    context = "\n\n".join([memory_texts[i] for i in I[0]])

    # ---- Adjust token range based on answer length ----
    if length == "balanced":
        max_tokens = 800
        instruction = "Write a clear, well-developed explanation (2–4 paragraphs)."
    elif length == "detailed":
        max_tokens = 1500
        instruction = "Write a detailed, structured report (about one page, with clarity and flow)."
    else:
        max_tokens = 250
        instruction = "Provide a concise, direct answer (2–5 sentences)."

    # ---- Construct prompt ----
    prompt = f"""
You are a precise and factual AI assistant.

Your goal is to answer the user's question based strictly on the provided context.
Do not add assumptions or outside information.
{instruction}

If the context does not contain enough information to answer fully, state that briefly.

Context:
{context}

User Question:
{question}
"""

    model_order = ["gpt-5", "gpt-4o", "gpt-4o-mini"]
    last_err = None
    for m in model_order:
        try:
            completion = client.chat.completions.create(
                model=m,
                messages=[
                    {"role": "system", "content": "You are factual, organized, and professional."},
                    {"role": "user", "content": prompt},
                ],
                max_tokens=max_tokens,
                temperature=0.3
            )
            answer = completion.choices[0].message.content
            return jsonify({
                "answer": answer,
                "model_used": m,
                "mode": length
            })
        except Exception as e:
            last_err = str(e)
            continue

    return jsonify({"answer": f"Error generating answer: {last_err}"}), 500


# -------------------------------------------------------------------
# ♻️ RESET MEMORY
# -------------------------------------------------------------------
@app.route("/reset", methods=["POST"])
def reset_memory():
    """Clear short-term memory."""
    global memory_vectors, memory_texts
    memory_vectors.clear()
    memory_texts.clear()
    return jsonify({"message": "♻️ Short-term memory cleared."})


# -------------------------------------------------------------------
# ✅ MAIN ENTRY
# -------------------------------------------------------------------
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000)
